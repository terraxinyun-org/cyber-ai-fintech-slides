"""
Generate PPTX slide deck for Intro to Cybersecurity & AI for FinTech.
Theory-focused — technical hands-on exercises are left for the lab.
Uses Terra新云 brand: dark background (#0a0a0f), accent indigo (#6366f1), Space Grotesk.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──
BG_PRIMARY = RGBColor(0x0A, 0x0A, 0x0F)
BG_SECONDARY = RGBColor(0x12, 0x12, 0x1A)
BG_CARD = RGBColor(0x1A, 0x1A, 0x28)
ACCENT = RGBColor(0x63, 0x66, 0xF1)
ACCENT_LIGHT = RGBColor(0x81, 0x8C, 0xF8)
ACCENT_DARK = RGBColor(0x4F, 0x46, 0xE5)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
MAGENTA = RGBColor(0xD9, 0x46, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_SECONDARY = RGBColor(0xB0, 0xB0, 0xC0)
TEXT_MUTED = RGBColor(0x80, 0x80, 0x90)
SUCCESS = RGBColor(0x28, 0xA7, 0x45)
WARNING = RGBColor(0xFF, 0xC1, 0x07)
ERROR = RGBColor(0xDC, 0x35, 0x45)
BORDER = RGBColor(0x2A, 0x2A, 0x3A)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ── Helper Functions ──

def set_bg(slide, color=BG_PRIMARY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    shape.rotation = 0
    # Adjust corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return p


def add_para(tf, text, size=18, color=TEXT_SECONDARY, bold=False, space_before=Pt(4), space_after=Pt(2), bullet=False, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    if bullet:
        p.bullet = True  # Not supported directly, use prefix
    return p


def add_bullet_list(tf, items, size=16, color=TEXT_SECONDARY, start_new=True, level=0, bold_prefix=False):
    """Add bullet items to a text frame. Items can be str or (bold_part, rest) tuples."""
    for item in items:
        p = tf.add_paragraph()
        if isinstance(item, tuple):
            # Bold prefix + normal text
            run1 = p.add_run()
            run1.text = item[0]
            run1.font.size = Pt(size)
            run1.font.color.rgb = ACCENT_LIGHT
            run1.font.bold = True
            run1.font.name = 'Calibri'
            run2 = p.add_run()
            run2.text = item[1]
            run2.font.size = Pt(size)
            run2.font.color.rgb = color
            run2.font.bold = False
            run2.font.name = 'Calibri'
        else:
            p.text = f"  {item}"
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.font.name = 'Calibri'
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        p.level = level
    return tf


def add_accent_line(slide, left, top, width):
    """Add a thin accent-colored line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def make_title_slide(title, subtitle="", meta=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_bg(slide)

    # Gradient accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Title
    tb = add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.5))
    set_text(tb.text_frame, title, size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    if subtitle:
        tb2 = add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10.3), Inches(0.8))
        set_text(tb2.text_frame, subtitle, size=22, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    if meta:
        tb3 = add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5))
        set_text(tb3.text_frame, meta, size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # Logo bottom right
    add_logo(slide)
    return slide


def make_section_slide(module_num, title, duration=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Module number
    tb1 = add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(0.6))
    set_text(tb1.text_frame, f"MODULE {module_num:02d}", size=18, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # Title
    tb2 = add_text_box(slide, Inches(1.5), Inches(3.1), Inches(10.3), Inches(1.2))
    set_text(tb2.text_frame, title, size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Duration
    if duration:
        tb3 = add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.5))
        set_text(tb3.text_frame, duration, size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    add_logo(slide)
    return slide


def make_content_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # Accent bar top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Title
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.7))
    set_text(tb.text_frame, title, size=30, color=ACCENT_LIGHT, bold=True)

    # Underline
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(2))

    add_logo(slide)
    return slide


def add_logo(slide):
    """Add Terra新云 text logo bottom-right."""
    tb = add_text_box(slide, Inches(10.5), Inches(6.9), Inches(2.5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Terra"
    run.font.size = Pt(12)
    run.font.color.rgb = ACCENT
    run.font.bold = True
    run.font.name = 'Calibri'
    run2 = p.add_run()
    run2.text = "新云"
    run2.font.size = Pt(12)
    run2.font.color.rgb = WHITE
    run2.font.bold = True
    run2.font.name = 'Calibri'
    p.alignment = PP_ALIGN.RIGHT


def add_card(slide, left, top, width, height, title_text, items, title_color=ACCENT_LIGHT, border_col=BORDER):
    """Add a rounded card with title and bullet items."""
    shape = add_shape(slide, left, top, width, height, border_color=border_col)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)

    set_text(tf, title_text, size=16, color=title_color, bold=True)
    for item in items:
        if isinstance(item, tuple):
            p = tf.add_paragraph()
            r1 = p.add_run()
            r1.text = item[0]
            r1.font.size = Pt(13)
            r1.font.color.rgb = ACCENT_LIGHT
            r1.font.bold = True
            r1.font.name = 'Calibri'
            r2 = p.add_run()
            r2.text = item[1]
            r2.font.size = Pt(13)
            r2.font.color.rgb = TEXT_SECONDARY
            r2.font.name = 'Calibri'
            p.space_before = Pt(3)
        else:
            p = tf.add_paragraph()
            p.text = f"  {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_SECONDARY
            p.font.name = 'Calibri'
            p.space_before = Pt(3)
    return shape


def add_highlight_box(slide, left, top, width, height, text, text_size=15):
    """Add accent-bordered highlight box."""
    # Left border accent
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(4), height)
    border.fill.solid()
    border.fill.fore_color.rgb = ACCENT
    border.line.fill.background()

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Pt(4), top, width - Pt(4), height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x25)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(10)
    set_text(tf, text, size=text_size, color=TEXT_SECONDARY)
    return shape


# ════════════════════════════════════════════════════════════════
#  SLIDES
# ════════════════════════════════════════════════════════════════

# ── TITLE SLIDE ──
make_title_slide(
    "Intro to Cybersecurity\n& AI for FinTech",
    "OpenCode Edition  |  Hands-On Lab",
    "2.5 Hours  •  5 Modules  •  AI-Guided  •  Singapore Context"
)

# ── AGENDA ──
slide = make_content_slide("Agenda")
items = [
    ("Module 01 — Threat Landscape", "  (20 min)  CIA triad, common attacks, AI, SG regulations"),
    ("Module 02 — Scan, Discover & Exploit", "  (45 min)  Hands-on vulnerability discovery"),
    ("Module 03 — Incident Response", "  (30 min)  SIEM investigation, brute force analysis"),
    ("Module 04 — Governance & Compliance", "  (25 min)  MAS TRM, stakeholder interview, gap analysis"),
    ("Module 05 — Wrap-Up", "  (20 min)  AI limitations, career paths, next steps"),
]
tb = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "", size=1)
for bold_part, rest in items:
    p = tf.add_paragraph()
    r1 = p.add_run()
    r1.text = bold_part
    r1.font.size = Pt(20)
    r1.font.color.rgb = WHITE
    r1.font.bold = True
    r1.font.name = 'Calibri'
    r2 = p.add_run()
    r2.text = rest
    r2.font.size = Pt(16)
    r2.font.color.rgb = TEXT_MUTED
    r2.font.name = 'Calibri'
    p.space_before = Pt(16)
    p.space_after = Pt(4)

add_highlight_box(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.7),
                  "Theory slides here  →  Technical hands-on exercises in the lab environment", 15)

# ── YOUR MISSION ──
slide = make_content_slide("Your Mission")
tb = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "You are a security consultant hired to assess VoltPay Digital,\na Singapore-based FinTech payment platform.", size=20, color=TEXT_SECONDARY)

# Stats row
stats = [
    ("85", "Employees"),
    ("SGD 180M", "Monthly Volume"),
    ("2,400", "Merchants"),
    ("Q2 2026", "MAS Inspection"),
]
for i, (num, label) in enumerate(stats):
    x = Inches(0.8 + i * 3.0)
    shape = add_shape(slide, x, Inches(2.6), Inches(2.7), Inches(1.3), border_color=RGBColor(0x3A, 0x3A, 0x5A))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(12)
    set_text(tf, num, size=28, color=ACCENT_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
    add_para(tf, label, size=12, color=TEXT_MUTED, space_before=Pt(2))
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

tb2 = add_text_box(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(1.5))
tf2 = tb2.text_frame
tf2.word_wrap = True
set_text(tf2, "", size=1)
add_para(tf2, "VoltPay must pass a MAS inspection and earn CSA Cyber Trust Mark certification", size=17, color=TEXT_SECONDARY)
add_para(tf2, "to secure a DBS partnership worth SGD 5–8M/year.", size=17, color=TEXT_SECONDARY, space_before=Pt(0))
add_para(tf2, "", size=6)
add_para(tf2, "Key People:", size=16, color=WHITE, bold=True, space_before=Pt(10))
add_para(tf2, "  CEO: Amanda Lim  |  CTO: Marcus Chen  |  IT Manager: Kevin Tan  |  CISO: Jason Toh (new)", size=14, color=TEXT_MUTED)

# ════════ MODULE 01 ════════

make_section_slide(1, "Threat Landscape", "20 minutes  •  Theory")

# ── CIA TRIAD ──
slide = make_content_slide("The CIA Triad")
tb = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
set_text(tb.text_frame, "The three pillars of information security:", size=17, color=TEXT_SECONDARY)

add_card(slide, Inches(0.8), Inches(1.8), Inches(3.7), Inches(3.5),
         "Confidentiality", [
             "Only authorized people can access data",
             "",
             ("VoltPay: ", "Merchant transaction data,"),
             "customer PII, and API keys",
             "must stay private.",
             "",
             ("Example: ", "IDOR vulnerability lets"),
             "anyone access any transaction",
         ])

add_card(slide, Inches(4.8), Inches(1.8), Inches(3.7), Inches(3.5),
         "Integrity", [
             "Data hasn't been tampered with",
             "",
             ("VoltPay: ", "Payment amounts, merchant"),
             "records, and audit logs must",
             "be accurate and unaltered.",
             "",
             ("Example: ", "SQL injection could modify"),
             "transaction amounts",
         ])

add_card(slide, Inches(8.8), Inches(1.8), Inches(3.7), Inches(3.5),
         "Availability", [
             "Systems work when needed",
             "",
             ("VoltPay: ", "Payment processing must"),
             "handle SGD 180M/month",
             "without downtime.",
             "",
             ("Example: ", "DDoS attack takes payment"),
             "gateway offline",
         ])

add_highlight_box(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.8),
                  "Every security decision maps back to CIA. Ask: \"Which pillar does this threat affect?\"", 15)

# ── COMMON ATTACKS ──
slide = make_content_slide("4 Common Cyber Attacks")

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.6), Inches(2.2),
         "Phishing", [
             "Fake emails/messages trick users into revealing credentials",
             "91% of cyberattacks start with phishing",
             ("VoltPay risk: ", "CFO gets fake 'urgent wire transfer' from 'CEO'"),
         ], border_col=RGBColor(0x44, 0x33, 0x55))

add_card(slide, Inches(6.7), Inches(1.3), Inches(5.6), Inches(2.2),
         "Ransomware", [
             "Malware encrypts all files, demands crypto payment",
             "Average ransom: USD 1.5M (2025)",
             ("VoltPay risk: ", "Payment DB encrypted → 2,400 merchants offline"),
         ], border_col=RGBColor(0x55, 0x33, 0x33))

add_card(slide, Inches(0.8), Inches(3.8), Inches(5.6), Inches(2.2),
         "DDoS (Distributed Denial of Service)", [
             "Overwhelm servers so legitimate users can't connect",
             "Attacks on financial services up 154% in 2025",
             ("VoltPay risk: ", "Payment gateway down during peak hours"),
         ], border_col=RGBColor(0x33, 0x44, 0x55))

add_card(slide, Inches(6.7), Inches(3.8), Inches(5.6), Inches(2.2),
         "Social Engineering", [
             "Manipulating people to bypass security controls",
             "Exploits trust, urgency, and authority",
             ("VoltPay risk: ", "Attacker calls helpdesk as 'Kevin from IT'"),
         ], border_col=RGBColor(0x44, 0x44, 0x33))

add_highlight_box(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6),
                  "Question: Which of these can technology alone prevent?  (Answer: None — humans are always a factor)", 14)

# ── AI IN SECURITY ──
slide = make_content_slide("AI in Cybersecurity — A Two-Sided Arms Race")

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.6), Inches(3.5),
         "Attackers Use AI", [
             "AI-generated phishing (no more typos/grammar errors)",
             "Deepfake voice calls impersonating executives",
             "Automated vulnerability scanning at massive scale",
             "AI-written malware that evades detection",
             "Faster reconnaissance and target profiling",
             "",
             "AI lowers the barrier — less skill needed to attack",
         ], title_color=ERROR, border_col=RGBColor(0x55, 0x25, 0x25))

add_card(slide, Inches(6.7), Inches(1.3), Inches(5.6), Inches(3.5),
         "Defenders Use AI", [
             "Anomaly detection in network traffic patterns",
             "Log analysis across millions of events per day",
             "Automated incident response playbooks",
             "Threat intelligence correlation",
             "AI coding assistants for secure development",
             "",
             "AI amplifies — analysts handle 10x more alerts",
         ], title_color=SUCCESS, border_col=RGBColor(0x20, 0x44, 0x30))

add_highlight_box(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.0),
                  "Today you'll use OpenCode — a free, open-source AI agent — as your security co-pilot.\nAI is a force multiplier, not a replacement for thinking.", 16)

# ── SINGAPORE REGULATIONS ──
slide = make_content_slide("Singapore Cyber Regulations")

add_card(slide, Inches(0.8), Inches(1.3), Inches(3.7), Inches(3.6),
         "MAS TRM", [
             "Monetary Authority of Singapore",
             "Technology Risk Management Guidelines",
             "",
             "Applies to: Banks, payment institutions",
             "",
             ("Key areas: ", ""),
             "Board governance (§3)",
             "Secure development (§6)",
             "Access control (§9)",
             "Data & infra security (§11)",
             "Cyber operations (§12)",
             "Vulnerability management (§13)",
         ], border_col=RGBColor(0x3A, 0x3A, 0x6A))

add_card(slide, Inches(4.8), Inches(1.3), Inches(3.7), Inches(3.6),
         "PDPA", [
             "Personal Data Protection Act",
             "",
             "Applies to: All organisations",
             "collecting personal data in SG",
             "",
             ("Consent: ", "Purpose limitation"),
             ("Breach: ", "Notify within 3 days"),
             ("Fines: ", "Up to SGD 1M or"),
             "10% of annual turnover",
             "",
             "VoltPay has 2,400 merchants'",
             "personal + financial data",
         ], border_col=RGBColor(0x3A, 0x5A, 0x3A))

add_card(slide, Inches(8.8), Inches(1.3), Inches(3.7), Inches(3.6),
         "CSA Cyber Trust", [
             "Cyber Security Agency of SG",
             "Voluntary certification mark",
             "",
             "5 domains, 22 areas assessed:",
             "Governance, protection,",
             "detection, response, awareness",
             "",
             ("VoltPay need: ", "Required by"),
             "DBS for the partnership",
             "(worth SGD 5–8M/year)",
             "",
             "Cert valid for 3 years",
         ], border_col=RGBColor(0x5A, 0x5A, 0x3A))

add_highlight_box(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.8),
                  "VoltPay must comply with all three — MAS TRM for payment licence, PDPA for merchant data, CSA for DBS deal.", 15)

# ════════ MODULE 02 ════════

make_section_slide(2, "Scan, Discover & Exploit", "45 minutes  •  Hands-On Lab")

# ── ATTACK METHODOLOGY ──
slide = make_content_slide("Attack Methodology")
tb = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
set_text(tb.text_frame, "Security testing follows the same structured approach used by professional penetration testers:", size=17, color=TEXT_SECONDARY)

# Attack flow boxes
steps = [
    ("1. Recon", "Gather information\nabout the target"),
    ("2. Scanning", "Find open ports\nand services"),
    ("3. Exploitation", "Test vulnerabilities\nfor real impact"),
    ("4. Post-Exploit", "Determine what\ndata is accessible"),
    ("5. Reporting", "Document findings\nand recommend fixes"),
]
for i, (title, desc) in enumerate(steps):
    x = Inches(0.8 + i * 2.5)
    shape = add_shape(slide, x, Inches(2.0), Inches(2.2), Inches(1.6), border_color=RGBColor(0x3A, 0x3A, 0x6A))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(10)
    set_text(tf, title, size=16, color=ACCENT_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
    add_para(tf, desc, size=12, color=TEXT_MUTED, space_before=Pt(8))
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

# Arrow connectors (simple right arrows between boxes)
for i in range(4):
    x = Inches(3.0 + i * 2.5)
    tb = add_text_box(slide, x, Inches(2.5), Inches(0.3), Inches(0.4))
    set_text(tb.text_frame, "→", size=20, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

tb2 = add_text_box(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(2.5))
tf2 = tb2.text_frame
tf2.word_wrap = True
set_text(tf2, "What Students Will Discover in the Lab", size=20, color=WHITE, bold=True)
add_para(tf2, "", size=6)
items2 = [
    ("Port Scanning (nmap) — ", "Which services are running and exposed?"),
    ("Vulnerability Scanning (nikto) — ", "What security misconfigurations exist?"),
    ("SQL Injection — ", "Can we extract data from the database?"),
    ("Cross-Site Scripting (XSS) — ", "Can we run code in users' browsers?"),
    ("Broken Access Control (IDOR) — ", "Can we access other users' data?"),
    ("Default Credentials — ", "Are factory passwords still in use?"),
]
for bold, rest in items2:
    p = tf2.add_paragraph()
    r1 = p.add_run()
    r1.text = bold
    r1.font.size = Pt(15)
    r1.font.color.rgb = ACCENT_LIGHT
    r1.font.bold = True
    r1.font.name = 'Calibri'
    r2 = p.add_run()
    r2.text = rest
    r2.font.size = Pt(15)
    r2.font.color.rgb = TEXT_SECONDARY
    r2.font.name = 'Calibri'
    p.space_before = Pt(4)

# ── OWASP TOP 10 ──
slide = make_content_slide("OWASP Top 10 (2021)")
tb = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
set_text(tb.text_frame, "The industry-standard ranking of the most critical web application security risks:", size=17, color=TEXT_SECONDARY)

owasp_items = [
    ("A01", "Broken Access Control", "Users can act outside their intended permissions", "CRITICAL"),
    ("A02", "Cryptographic Failures", "Sensitive data exposed due to weak/missing encryption", "HIGH"),
    ("A03", "Injection", "Untrusted data sent to an interpreter (SQL, XSS, etc.)", "CRITICAL"),
    ("A05", "Security Misconfiguration", "Insecure defaults, missing headers, verbose errors", "MEDIUM"),
    ("A07", "Auth Failures", "Weak passwords, credential stuffing, no rate limiting", "HIGH"),
    ("A08", "Software Integrity", "Code/data modified without verification", "MEDIUM"),
]

y_start = Inches(1.9)
for i, (code, name, desc, severity) in enumerate(owasp_items):
    y = y_start + Inches(i * 0.7)
    # Row background
    if i % 2 == 0:
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(0.65))
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = RGBColor(0x12, 0x12, 0x1F)
        row_bg.line.fill.background()

    # Code
    tb1 = add_text_box(slide, Inches(0.9), y + Pt(6), Inches(0.8), Inches(0.5))
    set_text(tb1.text_frame, code, size=14, color=ACCENT, bold=True)

    # Name
    tb2 = add_text_box(slide, Inches(1.8), y + Pt(6), Inches(2.8), Inches(0.5))
    set_text(tb2.text_frame, name, size=14, color=WHITE, bold=True)

    # Description
    tb3 = add_text_box(slide, Inches(4.8), y + Pt(6), Inches(6.0), Inches(0.5))
    set_text(tb3.text_frame, desc, size=13, color=TEXT_SECONDARY)

    # Severity badge
    sev_color = ERROR if severity == "CRITICAL" else (WARNING if severity == "HIGH" else ACCENT_LIGHT)
    tb4 = add_text_box(slide, Inches(11.2), y + Pt(6), Inches(1.2), Inches(0.5))
    set_text(tb4.text_frame, severity, size=11, color=sev_color, bold=True, alignment=PP_ALIGN.RIGHT)

add_highlight_box(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.6),
                  "Students will discover vulnerabilities mapped to A01, A03, A05, and A07 in the lab.", 14)

# ── WHAT STOLEN DATA IS WORTH ──
slide = make_content_slide("What Stolen Data Is Worth")
tb = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5))
set_text(tb.text_frame, "On the dark web, FinTech data commands premium prices:", size=17, color=TEXT_SECONDARY)

data_prices = [
    ("Credit card + CVV", "USD 5–35", "Singapore cards at premium"),
    ("Bank login credentials", "USD 50–200", "Higher for SG banks"),
    ("Full identity package (Fullz)", "USD 30–100", "Name, IC, address, DOB"),
    ("Corporate database dump", "USD 1,000–10,000+", "VoltPay = 2,400 merchants"),
]
y = Inches(1.9)
# Header
for col, (x, w, text) in enumerate([
    (Inches(0.9), Inches(3.5), "Data Type"),
    (Inches(4.5), Inches(2.2), "Dark Web Price"),
    (Inches(6.8), Inches(4.5), "Notes"),
]):
    tb = add_text_box(slide, x, y, w, Inches(0.4))
    set_text(tb.text_frame, text, size=13, color=ACCENT, bold=True)
y += Inches(0.45)

for dtype, price, note in data_prices:
    tb1 = add_text_box(slide, Inches(0.9), y, Inches(3.5), Inches(0.4))
    set_text(tb1.text_frame, dtype, size=14, color=WHITE)
    tb2 = add_text_box(slide, Inches(4.5), y, Inches(2.2), Inches(0.4))
    set_text(tb2.text_frame, price, size=14, color=WARNING, bold=True)
    tb3 = add_text_box(slide, Inches(6.8), y, Inches(4.5), Inches(0.4))
    set_text(tb3.text_frame, note, size=13, color=TEXT_MUTED)
    y += Inches(0.45)

add_card(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.2),
         "VoltPay Breach Impact", [
             ("MAS: ", "Fines, licence conditions, possible revocation"),
             ("PDPA: ", "Fines up to SGD 1M (10% of annual turnover)"),
             ("Business: ", "2,400 merchants lose trust → revenue collapse"),
             ("Partnership: ", "DBS deal (SGD 5–8M/year) cancelled"),
             ("Legal: ", "Breach notification costs, potential lawsuits"),
             "",
             "Singapore data commands premium due to higher account balances and purchasing power.",
         ], title_color=ERROR, border_col=RGBColor(0x55, 0x25, 0x25))

# ════════ MODULE 03 ════════

make_section_slide(3, "Incident Response with AI", "30 minutes  •  SIEM Investigation")

# ── NIST IR ──
slide = make_content_slide("NIST Incident Response Lifecycle")

phases = [
    ("1. Preparation", "Tools, people, processes ready\nbefore an incident occurs", SUCCESS),
    ("2. Detection\n& Analysis", "Identify and investigate\nthe security event", ACCENT_LIGHT),
    ("3. Containment\n& Eradication", "Stop the attack, remove\nattacker access", WARNING),
    ("4. Post-Incident\nActivity", "Document, report, improve\ndefenses", ERROR),
]
for i, (title, desc, color) in enumerate(phases):
    x = Inches(0.8 + i * 3.1)
    shape = add_shape(slide, x, Inches(1.4), Inches(2.8), Inches(2.0))
    shape.line.color.rgb = color
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(12)
    set_text(tf, title, size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_para(tf, desc, size=12, color=TEXT_MUTED, space_before=Pt(10))
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

tb = add_text_box(slide, Inches(0.8), Inches(3.8), Inches(11.7), Inches(3.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "What is a SIEM?", size=22, color=WHITE, bold=True)
add_para(tf, "", size=4)
add_para(tf, "Security Information and Event Management — the command center of any SOC.", size=17, color=TEXT_SECONDARY)
add_para(tf, "Collects logs from every system, normalizes them, and enables search + correlation.", size=17, color=TEXT_SECONDARY, space_before=Pt(2))
add_para(tf, "", size=8)
add_para(tf, "Our SIEM Stack:", size=17, color=WHITE, bold=True)
add_para(tf, "  Elasticsearch — stores and searches 10,000+ log entries", size=15, color=TEXT_SECONDARY)
add_para(tf, "  Kibana — visualizes and explores data across 5 log indices", size=15, color=TEXT_SECONDARY)
add_para(tf, "  Log sources: Authentication, Firewall, DNS, Process (Sysmon), Web requests", size=15, color=TEXT_SECONDARY)

# ── INVESTIGATION OVERVIEW ──
slide = make_content_slide("Brute Force Investigation")
tb = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
set_text(tb.text_frame, "In the lab, students investigate a brute force attack embedded in 10,000 log entries:", size=17, color=TEXT_SECONDARY)

steps_data = [
    ("Step 1: Overview", "Count auth successes vs failures — massive ratio of failures = something is wrong"),
    ("Step 2: Top Offenders", "Aggregate by source IP — one IP has hundreds of failed attempts"),
    ("Step 3: Identify Targets", "Which usernames were targeted? Is this credential stuffing or a targeted attack?"),
    ("Step 4: Check Success", "Did the attacker eventually get in? Search for successful login from attacker IP"),
    ("Step 5: Cross-Reference", "Check firewall + DNS logs for the same IP — what else did they do post-compromise?"),
    ("Step 6: AI Synthesis", "Build attack timeline, extract IOCs, map to MITRE ATT&CK T1110 (Brute Force)"),
]
y = Inches(1.9)
for step, desc in steps_data:
    tb1 = add_text_box(slide, Inches(0.9), y, Inches(2.5), Inches(0.45))
    set_text(tb1.text_frame, step, size=14, color=ACCENT_LIGHT, bold=True)
    tb2 = add_text_box(slide, Inches(3.5), y, Inches(8.8), Inches(0.45))
    set_text(tb2.text_frame, desc, size=14, color=TEXT_SECONDARY)
    y += Inches(0.55)

add_highlight_box(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.0),
                  "Key skill: CORRELATION — looking at the same event from multiple data sources.\nThe value of a SIEM is connecting the dots that no single log source can show you.", 15)

# ── IR REPORT ──
slide = make_content_slide("Incident Response Report")
tb = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
set_text(tb.text_frame, "Students draft a professional IR report with AI assistance:", size=17, color=TEXT_SECONDARY)

add_card(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.3),
         "Report Structure", [
             ("1. Executive Summary — ", "3 sentences for the CEO"),
             ("2. Timeline — ", "Chronological attack events"),
             ("3. IOCs — ", "Attacker IP, usernames, timestamps"),
             ("4. Root Cause — ", "Why the attack succeeded"),
             ("5. Impact — ", "What was compromised"),
             ("6. Recommendations — ", "Immediate + long-term fixes"),
             ("7. Lessons Learned — ", "Prevent recurrence"),
             "",
             "Writing for the audience:",
             "  Executive Summary → CEO",
             "  IOCs → SOC team",
             "  Recommendations → IT team",
         ])

add_card(slide, Inches(6.7), Inches(1.8), Inches(5.6), Inches(4.3),
         "IOCs to Extract", [
             "Attacker IP address(es)",
             "Targeted usernames",
             "Successful compromise timestamp",
             "Post-compromise network activity",
             "MITRE ATT&CK: T1110 (Brute Force)",
             "",
             ("Why this matters: ", ""),
             "IR reports are legal documents.",
             "They may be submitted to MAS,",
             "used in court proceedings, or",
             "shared with law enforcement.",
             "",
             "Accuracy and completeness matter.",
         ], border_col=RGBColor(0x3A, 0x3A, 0x6A))

# ════════ MODULE 04 ════════

make_section_slide(4, "Governance & Compliance", "25 minutes  •  Stakeholder Interview")

# ── WHAT IS GOVERNANCE ──
slide = make_content_slide("What is IT Governance?")
tb = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Think of it like building codes for technology.", size=22, color=WHITE, bold=True)
add_para(tf, "Just as HDB has safety requirements for buildings, MAS has security requirements for payment systems.", size=18, color=TEXT_SECONDARY, space_before=Pt(8))

add_card(slide, Inches(0.8), Inches(2.8), Inches(3.7), Inches(3.0),
         "MAS TRM — Key Sections", [
             ("§3 ", "Board governance"),
             ("§6 ", "Secure development"),
             ("§9 ", "Access control"),
             ("§11 ", "Data & infrastructure"),
             ("§12 ", "Cyber operations (SOC)"),
             ("§13 ", "Vulnerability management"),
         ], border_col=RGBColor(0x3A, 0x3A, 0x6A))

add_card(slide, Inches(4.8), Inches(2.8), Inches(3.7), Inches(3.0),
         "Why It Matters for VoltPay", [
             "Major Payment Institution licence",
             "MAS can impose conditions or revoke",
             "Q2 2026 inspection approaching",
             "PDPA fines up to SGD 1M",
             "DBS requires Cyber Trust Mark",
             "Insurance premiums based on compliance",
         ], border_col=RGBColor(0x3A, 0x5A, 0x3A))

add_card(slide, Inches(8.8), Inches(2.8), Inches(3.7), Inches(3.0),
         "Not Just Checkbox", [
             "Governance = culture, not just paperwork",
             "",
             "Good governance prevents incidents",
             "Bad governance causes incidents",
             "",
             ("Example: ", "Kevin's checklist is in his"),
             "head — if he leaves, there is",
             "no offboarding process.",
         ], border_col=RGBColor(0x5A, 0x3A, 0x3A))

# ── KEVIN TAN ──
slide = make_content_slide("Meet Kevin Tan — IT Manager")
tb = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5))
set_text(tb.text_frame, "Students interview Kevin (AI persona) to discover infrastructure gaps:", size=17, color=TEXT_SECONDARY)

add_card(slide, Inches(0.8), Inches(1.9), Inches(5.6), Inches(3.0),
         "Kevin's Profile", [
             "34 years old, NTU Computer Engineering",
             "5 years at Singtel NOC → VoltPay (2021)",
             "Built entire IT infrastructure from scratch",
             "1 junior sysadmin (6 months in)",
             "Manages 85 endpoints + network + physical security",
             "",
             "Speaks heavy Singlish, practical, resourceful,",
             "frustrated by budget constraints",
             "",
             "\"Aiyoh, you think I don't know meh?",
             "I raised this before but nobody listen...\"",
         ])

add_card(slide, Inches(6.7), Inches(1.9), Inches(5.6), Inches(3.0),
         "Interview Technique", [
             ("Surface answers ", "sound fine"),
             ("Dig deeper ", "to find real issues"),
             "",
             "Good follow-up questions:",
             "  \"When was that last reviewed?\"",
             "  \"What happens when someone leaves?\"",
             "  \"Who has access to the root account?\"",
             "",
             "Bad questions (get generic answers):",
             "  \"Tell me about your security\"",
             "  \"Is everything secure?\"",
         ], border_col=RGBColor(0x3A, 0x5A, 0x3A))

add_highlight_box(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.8),
                  "This teaches a critical soft skill: getting the REAL answers from stakeholders. Same technique used by auditors, consultants, and GRC professionals.", 14)

# ── GAP ANALYSIS ──
slide = make_content_slide("Gap Analysis — What Kevin Reveals")
tb = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
set_text(tb.text_frame, "Instructor reference — surface answers vs. reality after probing:", size=15, color=TEXT_MUTED)

gaps = [
    ("Endpoints", "85 machines, Intune, Defender", "12 Win10 EOL, 60/85 enrolled, 847 unreviewed alerts", "§11", "CRITICAL"),
    ("Network", "Separate guest/corp WiFi", "Same VLAN, AWS SSH open to 0.0.0.0/0", "§11", "CRITICAL"),
    ("Offboarding", "HR notifies, AD disabled", "Checklist in Kevin's head, ex-employees on GitHub", "§9", "HIGH"),
    ("Physical", "Card access, locked server room", "Key duplicated 5x, no camera, no temp monitoring", "§11", "MEDIUM"),
    ("Passwords", "Individual accounts", "AWS root in Google Doc, admin123 shared by all devs", "§9", "CRITICAL"),
]

# Header row
headers = [("Topic", 1.5), ("Surface Answer", 2.8), ("Reality (After Probing)", 4.0), ("MAS", 0.7), ("Risk", 1.0)]
y = Inches(1.8)
x = Inches(0.9)
for h_text, h_width in headers:
    tb = add_text_box(slide, x, y, Inches(h_width), Inches(0.4))
    set_text(tb.text_frame, h_text, size=12, color=ACCENT, bold=True)
    x += Inches(h_width + 0.15)

y += Inches(0.45)
for topic, surface, reality, mas, risk in gaps:
    x = Inches(0.9)
    vals = [(topic, 1.5, WHITE, True), (surface, 2.8, TEXT_MUTED, False), (reality, 4.0, TEXT_SECONDARY, False)]
    for text, w, color, bold in vals:
        tb = add_text_box(slide, x, y, Inches(w), Inches(0.55))
        set_text(tb.text_frame, text, size=12, color=color, bold=bold)
        x += Inches(w + 0.15)
    tb = add_text_box(slide, x, y, Inches(0.7), Inches(0.55))
    set_text(tb.text_frame, mas, size=12, color=ACCENT_LIGHT, bold=True)
    x += Inches(0.85)
    risk_color = ERROR if risk == "CRITICAL" else (WARNING if risk == "HIGH" else ACCENT_LIGHT)
    tb = add_text_box(slide, x, y, Inches(1.0), Inches(0.55))
    set_text(tb.text_frame, risk, size=11, color=risk_color, bold=True)
    y += Inches(0.55)

add_highlight_box(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.7),
                  "Students map each finding to MAS TRM sections and estimate business impact in SGD.", 14)

# ════════ MODULE 05 ════════

make_section_slide(5, "Wrap-Up & Reflection", "20 minutes")

# ── AI LIMITATIONS ──
slide = make_content_slide("AI — Strengths & Limitations")

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.6), Inches(3.0),
         "What AI Does Well", [
             "Explains complex concepts in plain language",
             "Analyzes tool output quickly",
             "Maps findings to frameworks (OWASP, MITRE)",
             "Drafts reports and documentation",
             "Suggests next investigative steps",
             "Simulates personas for practice (Kevin!)",
             "Helps write and review code",
         ], title_color=SUCCESS, border_col=RGBColor(0x20, 0x44, 0x30))

add_card(slide, Inches(6.7), Inches(1.3), Inches(5.6), Inches(3.0),
         "Where AI Falls Short", [
             ("Hallucination — ", "can fabricate facts confidently"),
             ("No verification — ", "can't check if systems are patched"),
             ("Over-confidence — ", "rarely says \"I don't know\""),
             ("Context limits — ", "loses track in long sessions"),
             ("Lacks judgment — ", "can't weigh business trade-offs"),
             ("No accountability — ", "you sign the report, not the AI"),
         ], title_color=ERROR, border_col=RGBColor(0x55, 0x25, 0x25))

add_highlight_box(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(1.0),
                  "AI is a force multiplier, not a replacement for thinking.\nThe person using the AI is still responsible for the output.", 18)

# ── CAREER PATHS ──
slide = make_content_slide("Career Paths in Cybersecurity")

add_card(slide, Inches(0.8), Inches(1.3), Inches(3.7), Inches(3.8),
         "Entry-Level Roles", [
             ("SOC Analyst — ", "Monitor alerts,"),
             "investigate incidents",
             "(Today's Module 03 skills)",
             "",
             ("GRC Analyst — ", "Compliance,"),
             "audits, risk assessment",
             "(Today's Module 04 skills)",
             "",
             ("Security Engineer — ", ""),
             "Build and maintain defenses",
         ], border_col=RGBColor(0x20, 0x44, 0x30))

add_card(slide, Inches(4.8), Inches(1.3), Inches(3.7), Inches(3.8),
         "Intermediate Roles", [
             ("Penetration Tester — ", ""),
             "Authorized hacking",
             "(Today's Module 02 skills)",
             "",
             ("Incident Responder — ", ""),
             "Forensics, malware analysis",
             "",
             ("Cloud Security — ", ""),
             "AWS/Azure/GCP hardening",
             "",
             ("Detection Engineer — ", ""),
             "Build SIEM rules and alerts",
         ], border_col=RGBColor(0x3A, 0x3A, 0x6A))

add_card(slide, Inches(8.8), Inches(1.3), Inches(3.7), Inches(3.8),
         "Certifications & Resources", [
             "Entry: CompTIA Security+",
             "Intermediate: CEH, CREST",
             "Advanced: OSCP, CISSP",
             "",
             ("Singapore: ", ""),
             "CSA SG Cyber Talent",
             "SG Cyber Youth",
             "DSTA CTF competitions",
             "NUS Greyhats",
             "",
             ("Practice: ", "TryHackMe, HackTheBox"),
         ], border_col=RGBColor(0x5A, 0x3A, 0x3A))

# ── ACHIEVEMENTS ──
slide = make_content_slide("Achievement Badges")
tb = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
set_text(tb.text_frame, "Students earn badges by completing milestones in the lab:", size=17, color=TEXT_SECONDARY)

badges = [
    ("First Blood", "Find first vulnerability"),
    ("SQL Sorcerer", "Execute SQL injection"),
    ("Script Kiddie\nNo More", "Explain WHY XSS works"),
    ("Master Key", "Access admin panel"),
    ("Digital Detective", "Investigate brute force"),
    ("Report Writer", "Complete IR report"),
    ("Governance Guru", "Find 3+ gaps from Kevin"),
    ("Full Stack\nDefender", "Complete all 5 modules"),
]
for i, (name, desc) in enumerate(badges):
    row = i // 4
    col = i % 4
    x = Inches(0.8 + col * 3.1)
    y = Inches(1.9 + row * 2.4)
    shape = add_shape(slide, x, y, Inches(2.8), Inches(2.0), border_color=RGBColor(0x3A, 0x3A, 0x6A))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(14)
    set_text(tf, name, size=16, color=ACCENT_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
    add_para(tf, desc, size=12, color=TEXT_MUTED, space_before=Pt(8))
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

add_highlight_box(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
                  "1–3 badges = Great start   •   4–5 = Strong foundation   •   6–8 = Outstanding", 14)

# ── CLOSING ──
make_title_slide(
    "Thank You",
    "You're now a cybersecurity defender.",
    "terraxinyun.com  •  josh@terraxinyun.com\n\nIntro to Cybersecurity & AI for FinTech  •  OpenCode Edition"
)


# ── SAVE ──
output_path = os.path.join(os.path.dirname(__file__), "Intro-to-Cybersecurity-AI-for-FinTech.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
